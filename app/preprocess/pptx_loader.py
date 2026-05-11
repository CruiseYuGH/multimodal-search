"""pptx → slide 文本 Chunk + slide 内嵌图 Chunk。"""
from __future__ import annotations
import io
import numpy as np

from app.core.types import Chunk, Modality, SourceType
from app.core.registry import register
from app.core.exceptions import PreprocessError
from app.preprocess.base import BasePreprocessor
from app.preprocess.text_loader import split_text
from app.utils.hashing import sha256_file, make_chunk_id
from app.utils.io import ensure_absolute
from config import settings


def _iter_shapes(shapes):
    """递归遍历 shape（含 group）。"""
    for s in shapes:
        if s.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_shapes(s.shapes)
        else:
            yield s


@register((".pptx",))
class PptxLoader(BasePreprocessor):
    SUPPORTED_EXT = (".pptx",)

    def load(self, file_path: str) -> list[Chunk]:
        path = ensure_absolute(file_path)
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from PIL import Image
        except ImportError as e:
            raise PreprocessError("缺少 python-pptx / Pillow 依赖") from e

        try:
            prs = Presentation(path)
        except Exception as e:
            raise PreprocessError(f"解析 pptx 失败: {path}: {e}") from e

        file_hash = sha256_file(path)
        max_len = settings.SPLITTER.split_length
        min_size = settings.PPTX_IMAGE_MIN_SIZE
        chunks: list[Chunk] = []
        idx = 0

        for slide_no, slide in enumerate(prs.slides, start=1):
            # ---- 文本 ----
            text_parts: list[str] = []
            image_idx_in_slide = 0
            for shape in _iter_shapes(slide.shapes):
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            text_parts.append(t)
                # 图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image
                        blob = img.blob
                        with Image.open(io.BytesIO(blob)) as pil:
                            pil = pil.convert("RGB")
                            w, h = pil.size
                            if w < min_size or h < min_size:
                                continue
                            arr = np.asarray(pil, dtype=np.uint8)
                    except Exception:
                        continue
                    chunks.append(Chunk(
                        chunk_id=make_chunk_id(path, idx),
                        source_path=path,
                        source_type=SourceType.PPTX,
                        modality=Modality.IMAGE,
                        file_hash=file_hash,
                        slide_no=slide_no,
                        image=arr,
                        extra={"image_idx": image_idx_in_slide, "width": w, "height": h},
                    ))
                    idx += 1
                    image_idx_in_slide += 1

            slide_text = "\n".join(text_parts).strip()
            if slide_text:
                if len(slide_text) <= max_len:
                    chunks.append(Chunk(
                        chunk_id=make_chunk_id(path, idx),
                        source_path=path,
                        source_type=SourceType.PPTX,
                        modality=Modality.TEXT,
                        file_hash=file_hash,
                        slide_no=slide_no,
                        text=slide_text,
                        preview=slide_text[:200],
                    ))
                    idx += 1
                else:
                    for ss, ee, seg in split_text(slide_text):
                        chunks.append(Chunk(
                            chunk_id=make_chunk_id(path, idx),
                            source_path=path,
                            source_type=SourceType.PPTX,
                            modality=Modality.TEXT,
                            file_hash=file_hash,
                            slide_no=slide_no,
                            text=seg,
                            preview=seg[:200],
                            extra={"char_start": ss, "char_end": ee},
                        ))
                        idx += 1
        return chunks
