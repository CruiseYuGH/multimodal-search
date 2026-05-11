"""B 阶段冒烟：在 /tmp/mm_smoke 下生成合成样例，跑全部 loader，打印 chunk 摘要。

不依赖 Milvus 与模型，只验证 preprocess 输出结构是否正确。
"""
from __future__ import annotations
import os
import sys
import json
from pathlib import Path

# 让脚本能直接 `python scripts/smoke_preprocess.py`
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.preprocess import (  # noqa: F401  触发注册
    text_loader, docx_loader, xlsx_loader, pptx_loader, image_loader,
)
from app.core.registry import get_preprocessor_for, list_registered
from app.core.types import Modality


SMOKE_DIR = Path("/tmp/mm_smoke")


def make_samples() -> dict[str, Path]:
    SMOKE_DIR.mkdir(parents=True, exist_ok=True)
    out: dict[str, Path] = {}

    # txt：构造一段长中文，确保会触发分块
    txt_path = SMOKE_DIR / "sample.txt"
    paragraph = "这是一段用于测试的中文文本。" * 80 + "\n\n第二段，包含分隔符。" + "另一句很长的内容。" * 50
    txt_path.write_text(paragraph, encoding="utf-8")
    out["txt"] = txt_path

    # docx
    from docx import Document
    doc = Document()
    doc.add_heading("测试文档", level=1)
    for i in range(5):
        doc.add_paragraph(f"第 {i+1} 段：这是 docx 测试内容，验证段落 index 与切分。" * (1 if i < 4 else 30))
    table = doc.add_table(rows=2, cols=3)
    table.rows[0].cells[0].text = "姓名"
    table.rows[0].cells[1].text = "年龄"
    table.rows[0].cells[2].text = "城市"
    table.rows[1].cells[0].text = "张三"
    table.rows[1].cells[1].text = "30"
    table.rows[1].cells[2].text = "上海"
    docx_path = SMOKE_DIR / "sample.docx"
    doc.save(docx_path)
    out["docx"] = docx_path

    # xlsx
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "销售"
    ws.append(["产品", "数量", "金额"])
    ws.append(["苹果", 10, 50])
    ws.append(["香蕉", 20, 40])
    ws2 = wb.create_sheet("库存")
    ws2.append(["SKU", "库存"])
    ws2.append(["A001", 100])
    xlsx_path = SMOKE_DIR / "sample.xlsx"
    wb.save(xlsx_path)
    out["xlsx"] = xlsx_path

    # png（合成纯色 + 文字）
    from PIL import Image, ImageDraw
    im = Image.new("RGB", (320, 240), (200, 50, 50))
    ImageDraw.Draw(im).rectangle([40, 40, 280, 200], fill=(50, 200, 50))
    png_path = SMOKE_DIR / "sample.png"
    im.save(png_path)
    out["image"] = png_path

    # pptx（含文本与一张内嵌图）
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "第一张幻灯片"
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    tb.text_frame.text = "这是第一张幻灯片的正文内容，用于验证 slide_no=1。"
    s1.shapes.add_picture(str(png_path), Inches(5), Inches(3), Inches(2), Inches(1.5))

    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "第二张幻灯片"
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    tb2.text_frame.text = "第二张幻灯片，验证多 slide 处理。" * 30  # 触发切分

    pptx_path = SMOKE_DIR / "sample.pptx"
    prs.save(pptx_path)
    out["pptx"] = pptx_path

    return out


def summarize(chunks):
    by_mod: dict[str, int] = {}
    for c in chunks:
        by_mod[c.modality.value] = by_mod.get(c.modality.value, 0) + 1
    return {"total": len(chunks), "by_modality": by_mod}


def main() -> int:
    print(f"[registry] {list_registered()}")
    samples = make_samples()
    print(f"[samples] {SMOKE_DIR}")
    for k, p in samples.items():
        print(f"  - {k}: {p}")

    rc = 0
    for kind, path in samples.items():
        try:
            loader = get_preprocessor_for(str(path))
            chunks = loader.load(str(path))
        except Exception as e:
            print(f"[FAIL] {kind} {path}: {type(e).__name__}: {e}")
            rc = 1
            continue
        s = summarize(chunks)
        print(f"[OK ] {kind:<5} {s}")
        # 抽样打印前 2 个 chunk 的关键字段
        for c in chunks[:2]:
            view = {
                "chunk_id": c.chunk_id,
                "modality": c.modality.value,
                "page": c.page, "slide_no": c.slide_no, "sheet_name": c.sheet_name,
                "preview": (c.preview or "")[:60],
                "image_shape": None if c.image is None else list(c.image.shape),
                "extra": c.extra,
            }
            print("       ", json.dumps(view, ensure_ascii=False))
    return rc


if __name__ == "__main__":
    sys.exit(main())
